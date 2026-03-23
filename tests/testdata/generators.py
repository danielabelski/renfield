"""Document format generators for test data.

Generates documents in PDF, DOCX, XLSX, PPTX, HTML, and MD formats.
TXT is NOT supported by Renfield's Docling document processor — use MD instead.

Dependencies: pip install fpdf2 python-docx openpyxl python-pptx
"""

import tempfile
from pathlib import Path

# Default output directory (overridable)
OUT_DIR = Path(tempfile.mkdtemp(prefix="renfield_testdata_"))


def set_output_dir(path: Path):
    """Override the output directory for generated files."""
    global OUT_DIR
    OUT_DIR = path
    OUT_DIR.mkdir(parents=True, exist_ok=True)


def create_pdf(filename: str, title: str, content: str) -> Path:
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    # fpdf2 built-in Helvetica is latin-1 only — replace Umlaute
    safe_title = (title.replace("€", "EUR").replace("ü", "ue").replace("ö", "oe")
                  .replace("ä", "ae").replace("Ü", "Ue").replace("Ö", "Oe")
                  .replace("Ä", "Ae").replace("ß", "ss"))
    pdf.cell(0, 10, safe_title, ln=True)
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    safe = (content.replace("€", "EUR").replace("ü", "ue").replace("ö", "oe")
            .replace("ä", "ae").replace("Ü", "Ue").replace("Ö", "Oe")
            .replace("Ä", "Ae").replace("ß", "ss"))
    pdf.multi_cell(0, 6, safe)
    path = OUT_DIR / filename
    pdf.output(str(path))
    return path


def create_docx(filename: str, title: str, content: str) -> Path:
    from docx import Document
    doc = Document()
    doc.add_heading(title, level=1)
    for para in content.split("\n\n"):
        doc.add_paragraph(para.strip())
    path = OUT_DIR / filename
    doc.save(str(path))
    return path


def create_md(filename: str, title: str, content: str) -> Path:
    path = OUT_DIR / filename
    path.write_text(f"# {title}\n\n{content}", encoding="utf-8")
    return path


def create_html(filename: str, title: str, content: str) -> Path:
    paragraphs = "\n".join(
        f"<p>{p.strip()}</p>" for p in content.split("\n\n") if p.strip()
    )
    html = f"""<!DOCTYPE html>
<html lang="de">
<head><meta charset="utf-8"><title>{title}</title></head>
<body><h1>{title}</h1>{paragraphs}</body>
</html>"""
    path = OUT_DIR / filename
    path.write_text(html, encoding="utf-8")
    return path


def create_xlsx(filename: str, title: str, content: str) -> Path:
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]
    for i, line in enumerate(content.split("\n"), 1):
        parts = line.split("|") if "|" in line else [line]
        for j, cell in enumerate(parts, 1):
            ws.cell(row=i, column=j, value=cell.strip())
    path = OUT_DIR / filename
    wb.save(str(path))
    return path


def create_pptx(filename: str, title: str, content: str) -> Path:
    from pptx import Presentation
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "TechNova GmbH"
    # Content slides
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    for i in range(0, len(paragraphs), 2):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Folie {i // 2 + 1}"
        body = slide.placeholders[1]
        body.text = paragraphs[i]
        if i + 1 < len(paragraphs):
            body.text += "\n\n" + paragraphs[i + 1]
    path = OUT_DIR / filename
    prs.save(str(path))
    return path


GENERATORS = {
    "pdf": create_pdf,
    "docx": create_docx,
    "md": create_md,
    "html": create_html,
    "xlsx": create_xlsx,
    "pptx": create_pptx,
}


def generate(filename: str, title: str, content: str) -> Path:
    """Generate a document file in the format determined by filename extension."""
    ext = filename.rsplit(".", 1)[-1].lower()
    gen = GENERATORS.get(ext)
    if not gen:
        raise ValueError(f"Unknown format: {ext}")
    return gen(filename, title, content)
